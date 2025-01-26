from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.dml.color import RGBColor
import xml.etree.ElementTree as ET
import re

def parse_bpmn_xml(xml_file):
    """BPMN XMLファイルを解析し、形状とコネクタのデータを抽出する。"""
    try:
        with open(xml_file, 'r', encoding='utf-8') as f:
            xml_content = f.read()
            # ダブルクォートのエスケープ処理
            xml_content = re.sub(r'=""([^"]*?)""', r'="\1"', xml_content)
            
        tree = ET.fromstring(xml_content)
    except ET.ParseError as e:
        print(f"XML Parse Error: {e}")
        return None, None
    except Exception as e:
        print(f"Error reading or processing XML file: {e}")
        return None, None

    # BPMN名前空間
    bpmn_ns = {"bpmn": "http://www.omg.org/spec/BPMN/20100524/MODEL"}
    
    shapes_data = {}
    connectors = []

    # プロセス要素を取得
    process = tree.find("bpmn:process", namespaces=bpmn_ns)
    if process is None:
        print("Error: No process element found in BPMN XML.")
        return None, None

    # 形状のデータを抽出
    for element in process.findall("*", namespaces=bpmn_ns):
        if element.tag.endswith("startEvent") or element.tag.endswith("task") or element.tag.endswith("exclusiveGateway") or element.tag.endswith("endEvent"):
            shape_id = element.get("id")
            shape_name = element.get("name")
            shape_type = element.tag.split("}")[1]  # タグから形状タイプを取得
            shapes_data[shape_id] = {"name": shape_name, "type": shape_type}

    # BPMNDiagram要素を取得
    bpmndiagram = tree.find("bpmndi:BPMNDiagram", namespaces={"bpmndi": "http://www.omg.org/spec/BPMN/20100524/DI"})
    if bpmndiagram is None:
        print("Error: No BPMNDiagram element found in BPMN XML.")
        return None, None

    bpmnplane = bpmndiagram.find("bpmndi:BPMNPlane", namespaces={"bpmndi": "http://www.omg.org/spec/BPMN/20100524/DI"})
    if bpmnplane is None:
        print("Error: No BPMNPlane element found in BPMN XML.")
        return None, None

    # 形状の位置とサイズを抽出
    for shape_di in bpmnplane.findall("bpmndi:BPMNShape", namespaces={"bpmndi": "http://www.omg.org/spec/BPMN/20100524/DI"}):
        shape_id = shape_di.get("bpmnElement")
        bounds = shape_di.find("dc:Bounds", namespaces={"dc": "http://www.omg.org/spec/DD/20100524/DC"})
        if bounds is not None:
            x = float(bounds.get("x")) / 100  # 座標を調整 (100で割ることで間隔を狭める)
            y = float(bounds.get("y")) / 100  # 座標を調整 (100で割ることで間隔を狭める)
            width = float(bounds.get("width")) / 100  # サイズを調整
            height = float(bounds.get("height")) / 100  # サイズを調整
            if shape_id in shapes_data:
                shapes_data[shape_id]["x"] = x
                shapes_data[shape_id]["y"] = y
                shapes_data[shape_id]["width"] = width
                shapes_data[shape_id]["height"] = height

    # コネクタのデータを抽出
    for sequence_flow in process.findall("bpmn:sequenceFlow", namespaces=bpmn_ns):
        source_ref = sequence_flow.get("sourceRef")
        target_ref = sequence_flow.get("targetRef")
        connectors.append((source_ref, target_ref))

    return shapes_data, connectors

def create_powerpoint(shapes_data, connectors, output_file):
    """抽出したデータに基づいてPowerPointを作成する。"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # 形状の描画
    shapes = {}
    for key, data in shapes_data.items():
        shape_type = data["type"]
        if shape_type == "startEvent" or shape_type == "endEvent":
            shape_type = MSO_SHAPE.OVAL
        elif shape_type == "task":
            shape_type = MSO_SHAPE.RECTANGLE
        elif shape_type == "exclusiveGateway":
            shape_type = MSO_SHAPE.DIAMOND
        else:
            continue  # 不明な形状タイプはスキップ

        # 「荷下ろし後の確認」のy座標を調整
        if key == "Task_3":
            gateway_y = shapes_data["ExclusiveGateway_1"]["y"]
            gateway_height = shapes_data["ExclusiveGateway_1"]["height"]
            data["y"] = gateway_y + gateway_height / 2 + 0.5  # ゲートウェイの下に配置 (間隔を半分にする)
        
        shape = slide.shapes.add_shape(
            shape_type,
            Inches(data["x"]), Inches(data["y"]), Inches(data["width"]), Inches(data["height"])
        )
        
        # テキストの設定
        text_frame = shape.text_frame
        text_frame.text = data["name"]
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.size = Inches(0.1)  # 文字サイズを小さく設定
            paragraph.font.color.rgb = RGBColor(0, 0, 0)  # 文字色を黒に設定
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # 図形の色を白黒に設定
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 白
        shape.line.color.rgb = RGBColor(0, 0, 0)  # 黒
        
        shapes[key] = shape

    # コネクタの描画
    for source, target in connectors:
        if source not in shapes or target not in shapes:
            continue  # 形状が存在しない場合はスキップ
        source_shape = shapes[source]
        target_shape = shapes[target]
        
        if source_shape.top != target_shape.top:
            connector_type = MSO_CONNECTOR.ELBOW
        else:
            connector_type = MSO_CONNECTOR.STRAIGHT
        
        connector = slide.shapes.add_connector(
            connector_type, 0, 0, 0, 0
        )
        
        if connector_type == MSO_CONNECTOR.ELBOW:
            # エルボーコネクタの接続位置を調整
            if source_shape.left < target_shape.left:
                connector.begin_connect(source_shape, 3)
                connector.end_connect(target_shape, 1)
            else:
                connector.begin_connect(source_shape, 1)
                connector.end_connect(target_shape, 3)
        else:
            connector.begin_connect(source_shape, 3)
            connector.end_connect(target_shape, 1)
        
        connector.line.width = Inches(0.02)
        connector.line.color.rgb = RGBColor(0, 0, 0)  # コネクタの色を黒に設定

    prs.save(output_file)
    print(f"PPTファイル '{output_file}' が作成されました。")

if __name__ == "__main__":
    xml_file = "C:\\Users\\toshimitsu_fujiki\\Downloads\\bpmn_workflow_fixed.xml"  # BPMN XMLファイルのパス
    output_file = "C:\\Users\\toshimitsu_fujiki\\Downloads\\bpmn_workflow_from_xml.pptx"  # 出力PowerPointファイルのパス

    shapes_data, connectors = parse_bpmn_xml(xml_file)
    if shapes_data and connectors:
        create_powerpoint(shapes_data, connectors, output_file)