# src/backend/api/projects.py
from fastapi import APIRouter, HTTPException, Path, Response, Body, Depends
from sqlalchemy.orm import Session
from pydantic import BaseModel
from datetime import datetime
import pandas as pd
import os
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.dml.color import RGBColor
import xml.etree.ElementTree as ET
import re
import logging
import io
import json
from typing import List, Optional
from database import get_db, Project

# ロギング設定
logging.basicConfig(level=logging.DEBUG)
router = APIRouter()

# Pydanticモデルの定義
class ProjectBase(BaseModel):
    id: Optional[int] = None  # ここをオプショナルに変更
    user_id: Optional[str] = None
    customer_name: Optional[str] = None
    issues: Optional[str] = None
    is_archived: Optional[bool] = None
    bpmn_xml: Optional[str] = None
    solution_requirements: Optional[str] = None
    stage: Optional[str] = "営業"
    category: Optional[str] = "プロジェクト"
    slack_channel_id: Optional[str] = None
    slack_tag: Optional[str] = None
    box_folder_path: Optional[str] = None
    schedule: Optional[str] = ""

class ProjectCreate(BaseModel):
    customer_name: str
    issues: str
    stage: str = "営業"
    category: str = "プロジェクト"

class StageUpdate(BaseModel):
    stage: str

class ProjectUpdate(ProjectBase):
    user_id: str = None
    customer_name: str = None
    issues: str = None
    stage: str = None
    category: str = None
    schedule: Optional[str] = None
    is_archived: Optional[bool] = None

class ProjectOut(ProjectBase):
    id: int
    is_archived: bool
    bpmn_xml: str = ""
    solution_requirements: str = ""
    slack_channel_id: str = ""
    slack_tag: str = ""
    box_folder_path: str = ""

    class Config:
        from_attributes = True

class ArchiveProjectUpdate(BaseModel):
    is_archived: bool

class SlackUpdate(ProjectBase):
    channel_id: str
    tag: str = ""

class FlowUpdate(BaseModel):
    bpmn_xml: str

class RequirementsUpdate(BaseModel):
    solution_requirements: str

PROJECTS_CSV = os.path.join(os.path.dirname(__file__), '../../../data/projects.csv')

def replace_none_with_empty(project):
    """プロジェクトオブジェクトの特定のフィールドが None なら空文字に置き換える"""
    fields = [
        "bpmn_xml",
        "slack_channel_id",
        "slack_tag",
        "solution_requirements",
        "box_folder_path",
        "schedule",
    ]
    for field in fields:
        if getattr(project, field) is None:
            setattr(project, field, "")

# プロジェクトCSVの読み込み
def read_projects() -> pd.DataFrame:
    if not os.path.exists(PROJECTS_CSV):
        df = pd.DataFrame(columns=['id', 'customer_name', 'issues', 'is_archived', 'bpmn_xml', 'solution_requirements', 'stage', 'category', 'slack_channel_id', 'slack_tag', 'box_folder_path', 'schedule'])
        df.to_csv(PROJECTS_CSV, index=False)
    df = pd.read_csv(PROJECTS_CSV, dtype={'id': int, 'customer_name': str, 'issues': str, 'is_archived': bool, 'bpmn_xml': str, 'solution_requirements': str, 'stage': str, 'category': str, 'slack_channel_id': str, 'slack_tag': str, 'box_folder_path': str, 'schedule': str})
    # 必要な列が存在しない場合は追加
    for col in ['slack_channel_id', 'slack_tag', 'box_folder_path', 'schedule']:
        if col not in df.columns or df[col].isnull().all():
            df[col] = ''
    # NaN を空文字列やデフォルト値に置き換える
    df['bpmn_xml'] = df['bpmn_xml'].fillna('')
    df['solution_requirements'] = df['solution_requirements'].fillna('')
    df['stage'] = df['stage'].fillna('営業')
    df['category'] = df['category'].fillna('プロジェクト')
    df['slack_channel_id'] = df['slack_channel_id'].fillna('').astype(str)
    df['slack_tag'] = df['slack_tag'].fillna('').astype(str)
    df['box_folder_path'] = df['box_folder_path'].fillna('').astype(str)
    df['schedule'] = df['schedule'].fillna('').astype(str)

    return df

# プロジェクトCSVへの書き込み
def write_projects(df: pd.DataFrame):
    df.to_csv(PROJECTS_CSV, index=False)

# プロジェクトの取得
@router.get("/", response_model=List[ProjectOut])
async def get_projects(db: Session = Depends(get_db)):
    projects = db.query(Project).all()
    for project in projects:
        replace_none_with_empty(project)
    return projects

# 新規プロジェクトの作成
@router.post("/", response_model=ProjectOut, status_code=201)
async def create_project(project: ProjectCreate, db: Session = Depends(get_db)):
    new_project = Project(
        user_id="user_888", # ユーザーIDをバックエンドで設定 (実際には認証などから取得)
        customer_name=project.customer_name,
        issues=project.issues,
        is_archived=False, # is_archived はデフォルトで False に設定
        stage=project.stage,
        category=project.category,
        bpmn_xml="",
        solution_requirements="",
        slack_channel_id="",
        slack_tag="",
        box_folder_path="",
    )
    logging.debug(f"作成するプロジェクトオブジェクト: {new_project.__dict__}")
    db.add(new_project)
    db.commit()
    db.refresh(new_project)
    return new_project

# プロジェクトの更新（顧客情報と課題）
@router.put("/{project_id}", response_model=ProjectOut)
async def update_project(project_id: int, project: ProjectUpdate, db: Session = Depends(get_db)):
    db_project = db.query(Project).filter(Project.id == project_id).first()
    if not db_project:
        raise HTTPException(status_code=404, detail="Project not found")

    if project.user_id:
        db_project.user_id = project.user_id
    if project.customer_name:
        db_project.customer_name = project.customer_name
    if project.issues:
        db_project.issues = project.issues
    if project.stage:
        db_project.stage = project.stage
    if project.category:
        db_project.category = project.category
    if project.schedule:
        db_project.schedule = project.schedule

    db.commit()
    db.refresh(db_project)
    return db_project

# アーカイブ状態の更新
@router.put("/{project_id}/archive", response_model=ProjectOut)
async def archive_project(
    project_id: int,
    project_update: ArchiveProjectUpdate = Body(...),
    db: Session = Depends(get_db)
):
    db_project = db.query(Project).filter(Project.id == project_id).first()
    if not db_project:
        raise HTTPException(status_code=404, detail="Project not found")

    if project_update.is_archived is not None:
        db_project.is_archived = project_update.is_archived

    db.commit()
    db.refresh(db_project)
    
    # ヘルパー関数で None を空文字に置換
    replace_none_with_empty(db_project)
    
    return db_project

@router.put("/{project_id}/stage", response_model=ProjectOut)
async def update_stage(
    project_id: int = Path(..., gt=0),
    stage_update: StageUpdate = Body(...),
    db: Session = Depends(get_db)
):
    db_project = db.query(Project).filter(Project.id == project_id).first()
    if not db_project:
        raise HTTPException(status_code=404, detail="Project not found")

    db_project.stage = stage_update.stage
    db.commit()
    db.refresh(db_project)
    
    # ヘルパー関数で None を空文字に置換
    replace_none_with_empty(db_project)
    
    return db_project

# 業務フローの更新
@router.put("/{project_id}/flow", response_model=ProjectOut)
async def update_flow(project_id: int = Path(..., gt=0), flow: FlowUpdate = None):
    df = read_projects()
    if project_id not in df['id'].values:
        raise HTTPException(status_code=404, detail="Project not found")
    index = df.index[df['id'] == project_id].tolist()[0]
    if flow.bpmn_xml:
        df.at[index, 'bpmn_xml'] = flow.bpmn_xml
    write_projects(df)
    updated_project = df.loc[index].to_dict()
    return ProjectOut(**updated_project)

@router.put("/{project_id}/requirements", response_model=ProjectOut)
async def update_requirements(project_id: int = Path(..., gt=0), req_update: RequirementsUpdate = None):
    df = read_projects()
    if project_id not in df['id'].values:
        raise HTTPException(status_code=404, detail="Project not found")
    index = df.index[df['id'] == project_id].tolist()[0]
    if req_update and req_update.solution_requirements is not None:
        df.at[index, 'solution_requirements'] = req_update.solution_requirements
    write_projects(df)
    updated_project = df.loc[index].to_dict()
    return ProjectOut(**updated_project)

def parse_bpmn_xml_content(xml_content):
    """BPMN XMLコンテンツを解析し、形状とコネクタのデータを抽出する。"""
    try:
        # ダブルクォートのエスケープ処理
        xml_content = re.sub(r'=""([^"]*?)""', r'="\1"', xml_content)
        tree = ET.fromstring(xml_content)
    except ET.ParseError as e:
        error_message = f"XML Parse Error: {e}"
        logging.error(error_message)
        raise HTTPException(status_code=400, detail=error_message)
    except Exception as e:
        error_message = f"Error reading or processing XML content: {e}"
        logging.error(error_message, exc_info=True)
        raise HTTPException(status_code=500, detail=error_message)

    # BPMN名前空間
    bpmn_ns = {"bpmn": "http://www.omg.org/spec/BPMN/20100524/MODEL"}

    shapes_data = {}
    connectors = []

    # プロセス要素を取得
    process = tree.find("bpmn:process", namespaces=bpmn_ns)
    if process is None:
        error_message = "Error: No process element found in BPMN XML."
        logging.error(error_message)
        raise HTTPException(status_code=400, detail=error_message)

    # 形状のデータを抽出
    for element in process.findall("*", namespaces=bpmn_ns):
        if element.tag.endswith("startEvent") or element.tag.endswith("task") or element.tag.endswith("exclusiveGateway") or element.tag.endswith("endEvent"):
            shape_id = element.get("id")
            shape_name = element.get("name")
            shape_type = element.tag.split("}")[1]  # タグから形状タイプを取得
            shapes_data[shape_id] = {"name": shape_name, "type": shape_type}

    # Gateway の outgoing SequenceFlow を抽出
    for gateway_element in process.findall("bpmn:exclusiveGateway", namespaces=bpmn_ns):
        gateway_id = gateway_element.get("id")
        outgoing_flows = []
        for sequence_flow in process.findall(f"bpmn:sequenceFlow[@sourceRef='{gateway_id}']", namespaces=bpmn_ns):
            outgoing_flows.append(sequence_flow.get("targetRef"))
        shapes_data[gateway_id]["outgoing"] = outgoing_flows

    # BPMNDiagram要素を取得
    bpmndiagram = tree.find("bpmndi:BPMNDiagram", namespaces={"bpmndi": "http://www.omg.org/spec/BPMN/20100524/DI"})
    if bpmndiagram is None:
        error_message = "Error: No BPMNDiagram element found in BPMN XML."
        logging.error(error_message)
        raise HTTPException(status_code=400, detail=error_message)

    bpmnplane = bpmndiagram.find("bpmndi:BPMNPlane", namespaces={"bpmndi": "http://www.omg.org/spec/BPMN/20100524/DI"})
    if bpmnplane is None:
        error_message = "Error: No BPMNPlane element found in BPMN XML."
        logging.error(error_message)
        raise HTTPException(status_code=400, detail=error_message)

    # 形状の位置とサイズを抽出
    for shape_di in bpmnplane.findall("bpmndi:BPMNShape", namespaces={"bpmndi": "http://www.omg.org/spec/BPMN/20100524/DI"}):
        shape_id = shape_di.get("bpmnElement")
        bounds = shape_di.find("dc:Bounds", namespaces={"dc": "http://www.omg.org/spec/DD/20100524/DC"})
        if bounds is not None:
            x = float(bounds.get("x")) / 100  # 座標を調整 (100で割ることで間隔を狭める)
            y = float(bounds.get("y")) / 100  # 座標を調整 (100で割ることで間隔を狭める)
            width = float(bounds.get("width")) / 100  # サイズを調整
            height = float(bounds.get("height")) / 100  # サイズを調整
            if shape_id in shapes_data:
                shapes_data[shape_id]["x"] = x
                shapes_data[shape_id]["y"] = y
                shapes_data[shape_id]["width"] = width
                shapes_data[shape_id]["height"] = height

    # コネクタのデータを抽出
    for sequence_flow in process.findall("bpmn:sequenceFlow", namespaces=bpmn_ns):
        source_ref = sequence_flow.get("sourceRef")
        target_ref = sequence_flow.get("targetRef")
        connectors.append((source_ref, target_ref))

    return shapes_data, connectors

def create_powerpoint_file(shapes_data, connectors):
    """抽出したデータに基づいてPowerPointを作成する (テンプレート使用、一時ファイルなし)。"""
    template_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), '../../../data/template.pptx')
    prs = Presentation(template_path)
    slide = prs.slides[2]

    y_offset = Inches(2)

    # 形状の描画
    shapes = {}
    for key, data in shapes_data.items():
        shape_type_str = data["type"] # 文字列の形状タイプ
        if shape_type_str == "startEvent" or shape_type_str == "endEvent":
            shape_type = MSO_SHAPE.OVAL
        elif shape_type_str == "task":
            shape_type = MSO_SHAPE.RECTANGLE
        elif shape_type_str == "exclusiveGateway":
            shape_type = MSO_SHAPE.DIAMOND
        else:
            continue  # 不明な形状タイプはスキップ

        # Gateway に続くオブジェクトの Y 座標を調整
        if data["type"] == "exclusiveGateway" and "outgoing" in data and data["outgoing"]:
            first_outgoing_task_id = data["outgoing"][0]
            if first_outgoing_task_id in shapes_data:
                gateway_y = data["y"]
                gateway_height = data["height"]
                shapes_data[first_outgoing_task_id]["y"] = gateway_y + gateway_height / 2 + 0.5  # ゲートウェイの下に配置 (間隔を半分にする)

        shape = slide.shapes.add_shape(
            shape_type,
            Inches(data["x"]), Inches(data["y"]) + y_offset, Inches(data["width"]), Inches(data["height"])
        )

        # テキストの設定
        text_frame = shape.text_frame
        text_frame.text = data["name"]
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.size = Inches(0.1)
            paragraph.font.name = "Meiryo UI"
            paragraph.font.color.rgb = RGBColor(0, 0, 0)
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        shape.line.color.rgb = RGBColor(0, 0, 0)

        shapes[key] = shape

    # コネクタの描画
    for source, target in connectors:
        if source not in shapes or target not in shapes:
            logging.warning(f"Connection source or target shape not found: {source} -> {target}")
            continue  # 形状が存在しない場合はスキップ
        source_shape = shapes[source]
        target_shape = shapes[target]

        if source_shape.top != target_shape.top:
            connector_type = MSO_CONNECTOR.ELBOW
        else:
            connector_type = MSO_CONNECTOR.STRAIGHT

        connector = slide.shapes.add_connector(
            connector_type, 0, 0, 0, 0
        )

        if connector_type == MSO_CONNECTOR.ELBOW:

            if source_shape.left < target_shape.left:
                connector.begin_connect(source_shape, 3)
                connector.end_connect(target_shape, 1)
            else:
                connector.begin_connect(source_shape, 1)
                connector.end_connect(target_shape, 3)
        else:
            connector.begin_connect(source_shape, 3)
            connector.end_connect(target_shape, 1)

        connector.line.width = Inches(0.02)
        connector.line.color.rgb = RGBColor(0, 0, 0)

    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io

# PowerPoint ファイルのエンドポイント (オブジェクト版)
@router.get("/{project_id}/powerpoint")
async def get_project_powerpoint(project_id: int = Path(..., gt=0)):
    df = read_projects()
    project_row = df[df['id'] == project_id]
    if project_row.empty:
        raise HTTPException(status_code=404, detail="Project not found")
    bpmn_xml = project_row.iloc[0]['bpmn_xml']
    if not bpmn_xml:
        raise HTTPException(status_code=400, detail="BPMN flow is not generated for this project.")

    pptx_file_path = None
    try:
        shapes_data, connectors = parse_bpmn_xml_content(bpmn_xml)
        if shapes_data and connectors:
            pptx_io = create_powerpoint_file(shapes_data, connectors)
            return Response(
                content=pptx_io.getvalue(), # メモリ上の PowerPoint ファイルのバイト列を取得
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                headers={"Content-Disposition": f"attachment; filename=project_{project_id}_flow_objects.pptx"} # headers
            )
        else:
            raise HTTPException(status_code=500, detail="Failed to parse BPMN XML and create shapes/connectors data.")
    except HTTPException as e:
        raise e
    except Exception as e:
        error_message = f"PowerPoint オブジェクト生成エラー: {str(e)}"
        logging.error(error_message, exc_info=True)
        raise HTTPException(status_code=500, detail=error_message)
    finally:
        pass

# 業務フローの削除
@router.delete("/{project_id}/flow", response_model=ProjectOut)
async def delete_flow(project_id: int = Path(..., gt=0)):
    df = read_projects()
    if project_id not in df['id'].values:
        raise HTTPException(status_code=404, detail="Project not found")
    index = df.index[df['id'] == project_id].tolist()[0]
    df.at[index, 'bpmn_xml'] = ""
    write_projects(df)
    updated_project = df.loc[index].to_dict()
    return ProjectOut(**updated_project)

# プロジェクトの削除
@router.delete("/{project_id}", status_code=204)
async def delete_project(project_id: int, db: Session = Depends(get_db)):
    db_project = db.query(Project).filter(Project.id == project_id).first()
    if not db_project:
        raise HTTPException(status_code=404, detail="Project not found")
    db.delete(db_project)
    db.commit()
    return {"detail": "Project deleted successfully"}

@router.put("/{project_id}/slack", response_model=ProjectOut)
async def connect_slack_channel(
    project_id: int = Path(..., gt=0),
    slack_data: SlackUpdate = Body(...)
):
    """
    プロジェクトにSlackチャンネルを連携し、タグも登録する
    """
    df = read_projects()
    if project_id not in df['id'].values:
        raise HTTPException(status_code=404, detail="Project not found")

    index = df.index[df['id'] == project_id][0]
    df.at[index, 'slack_channel_id'] = slack_data.channel_id
    df.at[index, 'slack_tag'] = slack_data.tag
    write_projects(df)
    updated_project = df.loc[index].to_dict()
    return ProjectOut(**updated_project)

@router.delete("/{project_id}/slack", response_model=ProjectOut)
async def disconnect_slack_channel(project_id: int = Path(..., gt=0)):
    """
    プロジェクトからSlackチャンネル連携を解除（チャンネルIDやタグをクリア）
    """
    df = read_projects()
    if project_id not in df['id'].values:
        raise HTTPException(status_code=404, detail="Project not found")

    index = df.index[df['id'] == project_id][0]
    df.at[index, 'slack_channel_id'] = ""
    df.at[index, 'slack_tag'] = ""
    write_projects(df)
    updated_project = df.loc[index].to_dict()
    return ProjectOut(**updated_project)