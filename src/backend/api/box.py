# src/backend/api/box.py
from fastapi import APIRouter, HTTPException, Depends
from fastapi.responses import FileResponse
from pydantic import BaseModel
from datetime import datetime, timezone
from typing import List, Optional
import os
import glob
import shutil
from api.projects import read_projects, write_projects
import pandas as pd
from sqlalchemy.orm import Session, declarative_base
from database import get_db, Project, UploadedFile
import fitz  # PyMuPDF
import json
import re
import requests
from pptx import Presentation
import pytz
try:
    import pytesseract
    from PIL import Image
except ImportError:
    # OCRを利用しない場合は、必要に応じてpytesseractとPillowをインストールしてください
    pytesseract = None

router = APIRouter()

# JSTタイムゾーンを定義
JST = pytz.timezone('Asia/Tokyo')

# Pydanticモデルの定義
class ProjectBase(BaseModel):
    user_id: str
    customer_name: str
    issues: Optional[str] = None
    is_archived: bool = False
    bpmn_xml: Optional[str] = None
    solution_requirements: Optional[str] = None
    stage: Optional[str] = None
    category: Optional[str] = None
    slack_channel_id: Optional[str] = None
    slack_tag: Optional[str] = None
    box_folder_path: Optional[str] = None
    schedule: Optional[datetime] = None

    class Config:
        orm_mode = True  # SQLAlchemyモデルをPydanticモデルに適合させる

class UploadedFileResponse(BaseModel):
    id: int
    source_name: str
    source_path: str
    project_id: int
    creation_date: datetime
    processed: bool
    processed_text: Optional[str]

    class Config:
        orm_mode = True

class BoxFileItem(BaseModel): # Renamed from UploadedFile to BoxFileItem
    filename: str
    filepath: str
    project_id: int

class LocalFileRequest(BaseModel):
    folder_path: str
    project_id: int

class BoxFolderLinkRequest(BaseModel):
    folder_path: str  # 相対パス

class BoxFileItem(BaseModel):
    filename: str
    filepath: str

class BaseDirectoryRequest(BaseModel):
    new_base_directory: str

CONFIG_FILE = "../../data/config.json"

# projectsテーブルからデータを取得する関数（PostgreSQLを使用）
def read_projects(db: Session) -> List[Project]:
    return db.query(Project).all()

# projectsテーブルにデータを保存する関数（PostgreSQLを使用）
def write_projects(db: Session, project_data: List[Project]):
    for project in project_data:
        db.add(project)
    db.commit()

# プロジェクトの取得エンドポイント
@router.get("/projects", response_model=List[ProjectBase])  # Pydanticモデルをレスポンスに指定
async def get_projects(db: Session = Depends(get_db)):
    try:
        projects = db.query(Project).all()
        return projects
    except Exception as e:
        raise HTTPException(status_code=500, detail="Error reading projects from database: " + str(e))

# プロジェクトの作成エンドポイント
@router.post("/projects", response_model=ProjectBase)
async def create_project(project: ProjectBase, db: Session = Depends(get_db)):
    try:
        project_data = Project(**project.dict())
        db.add(project_data)
        db.commit()
        db.refresh(project_data)
        return project_data
    except Exception as e:
        raise HTTPException(status_code=500, detail="Error writing project to database: " + str(e))

def read_config() -> dict:
    if not os.path.exists(CONFIG_FILE):
        # 初期設定として空のベースディレクトリを設定
        with open(CONFIG_FILE, 'w') as f:
            json.dump({"box_base_directory": ""}, f)
    with open(CONFIG_FILE, 'r') as f:
        return json.load(f)

def write_config(config: dict):
    with open(CONFIG_FILE, 'w') as f:
        json.dump(config, f, indent=4)

def boxnote_json_to_markdown(json_data: dict) -> str:
    if not json_data or not json_data.get('doc') or not json_data['doc'].get('content'):
        return ''

    def process_content(content: list) -> str:
        markdown = ''
        if not content:
            return markdown
        for node in content:
            if node['type'] == 'heading':
                markdown += f"{'#' * node['attrs']['level']} {process_inline_content(node.get('content', []))}\n"
            elif node['type'] == 'paragraph':
                markdown += f"{process_inline_content(node.get('content', []))}\n\n"
            elif node['type'] == 'bullet_list':
                markdown += process_list(node.get('content', []), '* ')
            elif node['type'] == 'ordered_list':
                markdown += process_list(node.get('content', []), '1. ')
            elif node['type'] == 'list_item':
                markdown += process_content(node.get('content', []))
            elif node['type'] == 'horizontal_rule':
                markdown += '---\n'
            elif node['type'] == 'text':
                markdown += process_text(node)
            elif node['type'] == 'hard_break':
                markdown += '\n'
            elif node['type'] == 'image':
                markdown += f"![{node['attrs']['alt'] or 'image'}]({node['attrs']['src'] or node['attrs']['boxSharedLink']})\n"
            elif node['type'] == 'embed':
                markdown += f"[{node['attrs']['title']}]({node['attrs']['url']})\n"
            elif node['type'] == 'doc':
                markdown += process_content(node.get('content', []))
            else:
                markdown += process_content(node.get('content', []))
        return markdown

    def process_inline_content(content: list) -> str:
        inline_markdown = ''
        if not content:
            return inline_markdown
        for node in content:
            if node['type'] == 'text':
                inline_markdown += process_text(node)
            elif node['type'] == 'hard_break':
                inline_markdown += '\n'
            elif node['type'] == 'image':
                inline_markdown += f"![{node['attrs']['alt'] or 'image'}]({node['attrs']['src'] or node['attrs']['boxSharedLink']})"
            elif node['type'] == 'embed':
                inline_markdown += f"[{node['attrs']['title']}]({node['attrs']['url']})"
            elif node.get('content'):
                inline_markdown += process_inline_content(node.get('content', []))
        return inline_markdown

    def process_text(node: dict) -> str:
        text = node.get('text', '')
        if node.get('marks'):
            for mark in node['marks']:
                if mark['type'] == 'strong':
                    text = f"**{text}**"
                elif mark['type'] == 'highlight':
                    text = f"<mark style='background-color:{mark['attrs'].get('color', '')}'>{text}</mark>"
        return text

    def process_list(content: list, prefix: str) -> str:
        list_markdown = ''
        index = 1
        for item in content:
            if item['type'] == 'list_item':
                current_prefix = f"{index}. " if prefix == '1. ' else prefix
                item_content = process_content(item['content'])
                item_content = item_content.replace('\n+', ' ').strip()
                list_markdown += f"{current_prefix}{item_content}\n"
                if prefix == '1. ':
                    index += 1
        list_markdown += '\n'
        return list_markdown

    return process_content(json_data['doc'].get('content', []))

# ファイルを一時的にローカルにダウンロードする関数
def download_file(file_url: str, local_path: str):
    response = requests.get(file_url, stream=True)
    if response.status_code == 200:
        with open(local_path, 'wb') as f:
            shutil.copyfileobj(response.raw, f)
    else:
        raise HTTPException(status_code=400, detail="ファイルのダウンロードに失敗しました。")

def extract_text_from_pdf(file_path: str) -> str:
    """
    PDFからテキストを抽出する関数
      - まず PyMuPDF でテキスト抽出を試みる。
      - テキストが空の場合は、OCRを用いて画像から文字認識を行う（pytesseract 必要）。
      - 抽出結果からハイフン改行や不要な改行を除去する。
      - 各ページの冒頭にページ番号「Page N:」を付与し、ページ間は2改行で区切る。
    """
    extracted_text = ""
    try:
        with fitz.open(file_path) as doc:
            for page_num, page in enumerate(doc, start=1):
                try:
                    # まずテキスト抽出を試行
                    page_text = page.get_text("text")
                except Exception as e:
                    page_text = ""
                
                # テキストが十分でない場合、OCRにフォールバック
                if not page_text.strip() and pytesseract is not None:
                    try:
                        pix = page.get_pixmap()
                        img_data = pix.tobytes()  # 画像データをバイト列として取得
                        img = Image.open(io.BytesIO(img_data))
                        page_text = pytesseract.image_to_string(img)
                    except Exception as ocr_e:
                        # OCR処理でもエラーの場合は、空文字のまま次へ
                        page_text = ""
                
                # ノイズ除去・ハイフンや改行の処理
                page_text = remove_pdf_noise(page_text)
                
                # ページ番号を先頭に付与
                page_text = f"Page {page_num}:\n" + page_text.strip()
                
                # ページごとに2改行で区切って連結
                extracted_text += page_text + "\n\n"
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PDFの処理に失敗しました: {e}")
    
    return extracted_text

def remove_pdf_noise(text: str) -> str:
    """
    PDF固有のノイズを除去し、ハイフン改行や不要な改行を整形する関数
      - "trailer" や "%%EOF" の除去
      - 行末のハイフンと改行（例: "ex-\ntraction"）を単語として連結
      - シングル改行はスペースに置換し、段落（2連続改行）は維持する
    """
    # 不要な文字列の除去
    text = text.replace("trailer", "").replace("%%EOF", "").strip()
    # ハイフンで改行されている箇所を結合（例："ex-\ntraction" → "extraction"）
    text = re.sub(r'-\n', '', text)
    # シングル改行（改行記号の前後が改行でない場合）をスペースに置換
    text = re.sub(r'(?<!\n)\n(?!\n)', ' ', text)
    return text

# ファイルパスをUTF-8でデコードする
def safe_file_path(file_path: str) -> str:
    try:
        return file_path.encode('utf-8').decode('utf-8')
    except UnicodeDecodeError:
        # エンコードできなかった場合、エラーを発生させる
        raise HTTPException(status_code=400, detail=f"無効なファイルパス: {file_path}")

# テキスト抽出処理を行うAPIエンドポイント
@router.get("/extract-text-from-file/{project_id}/{filename}")
async def extract_text_from_file(project_id: int, filename: str, db: Session = Depends(get_db)):
    config = read_config()
    base_directory = config.get('box_base_directory', '')

    if not base_directory:
        raise HTTPException(status_code=400, detail="ベースディレクトリが設定されていません。")

    # プロジェクトの存在チェック
    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    relative_path = project.box_folder_path

    # ファイルパスをUTF-8でデコード
    file_path = os.path.join(base_directory, relative_path, filename)
    file_path = safe_file_path(file_path)
    print(f"File path: {file_path}")

    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="File not found")

    # キャッシュの確認
    uploaded_file = db.query(UploadedFile).filter(
        UploadedFile.project_id == project_id,
        UploadedFile.source_name == filename
    ).first()

    if uploaded_file and uploaded_file.processed and uploaded_file.processed_text:
        print("キャッシュヒット")
        return {"text": uploaded_file.processed_text}

    print("キャッシュミス")

    file_extension = filename.split('.')[-1].lower()
    extracted_text = ""

    if file_extension == 'pdf':
        extracted_text = extract_text_from_pdf(file_path)
    elif file_extension == 'boxnote':
        with open(file_path, 'r', encoding='utf-8') as file:
            boxnote_json = json.load(file)
            extracted_text = boxnote_json_to_markdown(boxnote_json)
    # PowerPoint（ppt, pptx）の場合
    elif file_extension in ['ppt', 'pptx']:
        try:
            from pptx.enum.shapes import PP_PLACEHOLDER
            prs = Presentation(file_path)
            slide_texts = []
            for i, slide in enumerate(prs.slides):
                # シェイプを上から順、かつ左から順にソート（テキストが空でないもの）
                sorted_shapes = sorted(
                    [shape for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()],
                    key=lambda s: (s.top, s.left)
                )
                
                title_text = ""
                subtitle_text = ""
                content_parts = []
                
                # プレースホルダーや名称で判定して、タイトル、サブタイトル／メッセージラインを抽出
                for shape in sorted_shapes:
                    # プレースホルダー判定
                    if shape.is_placeholder:
                        ph_type = shape.placeholder_format.type
                        if ph_type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]:
                            if not title_text:
                                title_text = shape.text.strip()
                            continue
                        elif ph_type == PP_PLACEHOLDER.SUBTITLE:
                            if not subtitle_text:
                                subtitle_text = shape.text.strip()
                            continue
                    # シェイプ名によるヒューリスティック判定
                    name_lower = shape.name.lower() if hasattr(shape, "name") else ""
                    if "title" in name_lower and not title_text:
                        title_text = shape.text.strip()
                        continue
                    if "subtitle" in name_lower and not subtitle_text:
                        subtitle_text = shape.text.strip()
                        continue
                    # 上記に該当しないものは、コンテンツとして扱う
                    content_parts.append(shape.text.strip())
                
                # もしタイトルもサブタイトルも取得できなかった場合は、単純にソート結果の全テキストを連結して出力する（フォールバック）
                if not title_text and not subtitle_text:
                    fallback_text = "\n".join([shape.text.strip() for shape in sorted_shapes])
                    slide_texts.append(fallback_text)
                else:
                    # 表紙か通常スライドかの判断（ヒューリスティック例）
                    # → タイトルとサブタイトルがあり、かつコンテンツがほぼ無い、または最初のスライドなら表紙とする
                    if title_text and subtitle_text and (len(content_parts) == 0 or i == 0):
                        slide_text = f"{title_text}\n{subtitle_text}"
                    else:
                        page_number = i + 1  # ページ番号（1オリジン）
                        lines = []
                        if title_text:
                            lines.append(f"{page_number} {title_text}")
                        if subtitle_text:
                            lines.append(subtitle_text)
                        if content_parts:
                            lines.extend(content_parts)
                        slide_text = "\n".join(lines)
                    slide_texts.append(slide_text)
            extracted_text = "\n\n".join(slide_texts)
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"PowerPointファイルの処理に失敗しました: {e}")
    else:
        with open(file_path, 'r', encoding='utf-8') as file:
            extracted_text = file.read()

    # キャッシュの保存または更新
    if uploaded_file:
        uploaded_file.processed = True
        uploaded_file.processed_text = extracted_text
        uploaded_file.creation_date = datetime.now(JST)
    else:
        # ファイルがDBに存在しない場合は、processed=Trueとprocessed_textを設定して新規作成 (通常はlist-local-filesで事前登録されるはず)
        new_file = UploadedFile(
            source_name=filename,
            source_path=file_path,  # フルパスを保存
            project_id=project_id,
            creation_date=datetime.now(JST),
            processed=True,
            processed_text=extracted_text
        )
        db.add(new_file)
    db.commit()

    return {"text": extracted_text}

# グローバルベースディレクトリの取得エンドポイント
@router.get("/base-directory", response_model=dict)
async def get_base_directory():
    config = read_config()
    return {"box_base_directory": config.get("box_base_directory", "")}

# グローバルベースディレクトリの設定エンドポイント
@router.put("/base-directory", response_model=dict)
async def set_base_directory(req: BaseDirectoryRequest):
    new_base_directory = req.new_base_directory

    # バリデーション: 絶対パスであること（UnixとWindows両方に対応）
    if not (new_base_directory.startswith('/') or re.match(r'^[A-Za-z]:/', new_base_directory)):
        raise HTTPException(status_code=400, detail="ベースディレクトリは絶対パスで指定してください。")

    # パスの存在確認・作成
    if not os.path.exists(new_base_directory):
        try:
            os.makedirs(new_base_directory, exist_ok=True)
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"ベースディレクトリの作成に失敗しました: {e}")

    config = read_config()
    config['box_base_directory'] = new_base_directory
    write_config(config)
    return {"box_base_directory": new_base_directory, "message": "ベースディレクトリが更新されました。"}

# プロジェクトにBoxフォルダの相対パスを紐づけるエンドポイント
@router.put("/connect/{project_id}", response_model=dict)
async def connect_box_folder(project_id: int, req: BoxFolderLinkRequest):
    config = read_config()
    base_directory = config.get('box_base_directory', '')
    if not base_directory:
        raise HTTPException(status_code=400, detail="ベースディレクトリが設定されていません。")
    relative_path = req.folder_path
    full_path = os.path.join(base_directory, relative_path)

    # パスのバリデーション
    if not relative_path:
        raise HTTPException(status_code=400, detail="相対パスが空です。")
    if not os.path.isabs(base_directory):
        raise HTTPException(status_code=400, detail="ベースディレクトリが絶対パスではありません。")
    # パスの正規化とベースディレクトリの確認
    normalized_full_path = os.path.normpath(full_path)
    if not normalized_full_path.startswith(os.path.normpath(base_directory)):
        raise HTTPException(status_code=400, detail="相対パスがベースディレクトリを超えています。")

    # ディレクトリの存在確認・作成
    if not os.path.exists(full_path):
        try:
            os.makedirs(full_path, exist_ok=True)
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Boxフォルダの作成に失敗しました: {e}")

    df = read_projects()
    if project_id not in df['id'].values:
        raise HTTPException(status_code=404, detail="Project not found")

    index = df.index[df['id'] == project_id][0]
    df.at[index, 'box_folder_path'] = relative_path
    write_projects(df)
    updated_project = df.loc[index].to_dict()
    return {
        "project_id": project_id,
        "box_folder_path": relative_path,
        "message": "Boxフォルダがプロジェクトに連携されました。",
        "project": updated_project
    }

# プロジェクトからBoxフォルダの連携を解除するエンドポイント
@router.delete("/disconnect/{project_id}", response_model=dict)
async def disconnect_box_folder(project_id: int):
    """
    プロジェクトからBoxフォルダの相対パスの紐づけを解除。
    CSVの box_folder_path をクリアする処理。
    """
    df = read_projects()
    if project_id not in df['id'].values:
        raise HTTPException(status_code=404, detail="Project not found")

    index = df.index[df['id'] == project_id][0]
    df.at[index, 'box_folder_path'] = ""
    write_projects(df)
    updated_project = df.loc[index].to_dict()
    return {
        "project_id": project_id,
        "message": "Boxフォルダがプロジェクトから切断されました。",
        "project": updated_project
    }

# ローカルファイルの一覧取得エンドポイント
@router.post("/list-local-files", response_model=List[UploadedFileResponse])
async def list_local_files(req: LocalFileRequest, db: Session = Depends(get_db)):
    folder_path = req.folder_path
    project_id = req.project_id
    if not os.path.isdir(folder_path):
        raise HTTPException(status_code=400, detail="指定されたパスはフォルダではありません。")

    files = glob.glob(os.path.join(folder_path, "*"))
    uploaded_files = []
    db_changed = False # データベースに変更があったかどうかを追跡するフラグ

    for file_path in files:
        if os.path.isfile(file_path):
            filename = os.path.basename(file_path)
            source_name = filename
            source_path = file_path
            file_last_modified = datetime.fromtimestamp(os.path.getmtime(file_path))

            # 同じsource_name と project_id でファイルを探す (source_path はユニークキーに含まれない)
            existing_file = db.query(UploadedFile).filter(
                UploadedFile.source_name == source_name,
                UploadedFile.project_id == project_id
            ).first()

            if existing_file:
                # 既存のレコードがあれば、ファイルの最終更新日時を比較
                if existing_file.creation_date < file_last_modified:
                    print(f"ファイルが更新されました: {filename}") # デバッグログ
                    existing_file.processed = False
                    existing_file.processed_text = None
                    existing_file.creation_date = file_last_modified # creation_dateを更新 (最終処理日時をファイル最終更新日時に更新)
                    db_changed = True # 更新があったのでフラグを立てる
                else:
                    print(f"ファイルは更新されていません: {filename}") # デバッグログ
                    pass # 更新不要
            else:
                # 新規レコードの場合のみ挿入
                print(f"新規ファイル: {filename}") # デバッグログ
                new_file = UploadedFile(
                    source_name=source_name,
                    source_path=source_path,
                    project_id=project_id,
                    creation_date=file_last_modified, # 新規作成時はファイル最終更新日時をcreation_dateに設定
                    processed=False, # 新規ファイルは未処理
                    processed_text=None
                )
                db.add(new_file)
                existing_file = new_file # 新規ファイルもuploaded_filesに追加するためにexisting_fileに代入
                db_changed = True # 挿入があったのでフラグを立てる

            if existing_file:
                uploaded_files.append(existing_file)

    # データベースに変更があった場合のみコミット
    if db_changed:
        db.commit()
        for file in uploaded_files: # refresh は commit 後に行う
            db.refresh(file)

    return uploaded_files

# 全プロジェクトのファイル一覧取得エンドポイント
@router.get("/files", response_model=List[UploadedFileResponse])
async def get_all_files(project_id: int, db: Session = Depends(get_db)):
    """
    すべてのプロジェクトのファイル一覧を取得する。
    """
    files_from_db = db.query(UploadedFile).filter(UploadedFile.project_id == project_id).all() # データベースから UploadedFile モデルのリストを取得
    return files_from_db

# 特定プロジェクトのファイル一覧取得エンドポイント
@router.get("/files/{project_id}", response_model=List[UploadedFileResponse])
async def get_files_by_project(project_id: int, db: Session = Depends(get_db)):
    """
    特定のプロジェクトのファイル一覧を取得する。
    """
    config = read_config()
    base_directory = config.get('box_base_directory', '')
    if not base_directory:
        raise HTTPException(status_code=400, detail="ベースディレクトリが設定されていません。")
    
    db_files = db.query(UploadedFile).filter(UploadedFile.project_id == project_id).all()

    # SQLAlchemy モデル (database.UploadedFile) を Pydantic モデル (api.box.UploadedFileResponse) に変換
    response_files = [UploadedFileResponse.from_orm(db_file) for db_file in db_files]

    return response_files

# # ファイルアップロードエンドポイント
# @router.post("/upload-file/{project_id}", response_model=UploadedFile)
# async def upload_file(project_id: int, file: UploadFile = File(...)):
#     """
#     ファイルをアップロードする。
#     """
#     config = read_config()
#     base_directory = config.get('box_base_directory', '')
#     if not base_directory:
#         raise HTTPException(status_code=400, detail="ベースディレクトリが設定されていません。")
    
#     df = read_projects()
#     if project_id not in df['id'].values:
#         raise HTTPException(status_code=404, detail="Project not found")
    
#     project = df.loc[df['id'] == project_id].iloc[0]
#     relative_path = project.get('box_folder_path', '')
#     project_folder = os.path.join(base_directory, relative_path)
    
#     os.makedirs(project_folder, exist_ok=True)
#     file_path = os.path.join(project_folder, file.filename)
#     with open(file_path, "wb") as buffer:
#         shutil.copyfileobj(file.file, buffer)
#     return UploadedFile(filename=file.filename, filepath=file_path, project_id=project_id)

# # ファイル削除エンドポイント
# @router.delete("/delete-file/{project_id}/{filename}")
# async def delete_file(project_id: int, filename: str):
    """
    ファイルを削除する。
    """
    config = read_config()
    base_directory = config.get('box_base_directory', '')
    if not base_directory:
        raise HTTPException(status_code=400, detail="ベースディレクトリが設定されていません。")
    
    df = read_projects()
    if project_id not in df['id'].values:
        raise HTTPException(status_code=404, detail="Project not found")
    
    project = df.loc[df['id'] == project_id].iloc[0]
    relative_path = project.get('box_folder_path', '')
    file_path = os.path.join(base_directory, relative_path, filename)
    
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="File not found")
    os.remove(file_path)
    return {"message": "File deleted"}