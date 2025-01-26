# backend/api/proposals.py
from fastapi import APIRouter, HTTPException, Path, Response
from pydantic import BaseModel
from database import read_csv, write_csv
from llm_service import generate_proposal
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
import os
import io
from .projects import read_projects, PROJECTS_CSV # projects.py から関数と定数を import
from .projects import parse_bpmn_xml_content, create_powerpoint_file as create_bpmn_pptx # projects.py の関数を import (名前衝突を避けるため別名で import)
import logging
from openai import OpenAI
from dotenv import load_dotenv
import httpx
import json

load_dotenv()

# ログの設定
logging.basicConfig(level=logging.DEBUG)

# カスタム証明書のパス（自己署名証明書を指定）
CUSTOM_CERT_PATH = "C:\\Users\\toshimitsu_fujiki\\OpenAI_Cato_Networks_CA.crt"

# httpxクライアントを構築
httpx_client = httpx.Client(
    verify=CUSTOM_CERT_PATH  # 自己署名証明書を指定
)

# OpenAIクライアントにカスタムhttpxクライアントを渡す
api_key = os.getenv("OPENAI_API_KEY")
if not api_key:
    raise ValueError("APIキーが設定されていません。")

client = OpenAI(api_key=api_key, http_client=httpx_client)

logging.basicConfig(level=logging.DEBUG)

router = APIRouter()

TEMPLATE_PPTX_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), '../../../data/template.pptx') # テンプレートファイルパス

class ProposalCreate(BaseModel):
    customer_info: str
    project_info: str
    solution_id: str

def generate_slide2_content(customer_info: str, issues: str, solution_requirements: str):
    """LLMとFunction Callingを使ってスライド2枚目のコンテンツを生成する。"""

    functions = [
        {
            "name": "summarize_issues_and_solutions",
            "description": "顧客情報、課題、ソリューション要件から提案書スライド2枚目のコンテンツを生成",
            "parameters": {
                "type": "object",
                "properties": {
                    "current_issues": {
                        "type": "string",
                        "description": "顧客の現状の課題を箇条書きで要約",
                    },
                    "solution_direction": {
                        "type": "string",
                        "description": "ソリューション要件から考えられる解決の方向性を箇条書きで要約",
                    },
                    "message_line": {
                        "type": "string",
                        "description": "スライド2枚目で最も伝えたいメッセージを1文で記述",
                    },
                },
                "required": ["current_issues", "solution_direction", "message_line"],
            },
        }
    ]

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "system", "content": "あなたはPowerPoint提案資料作成のエキスパートです。与えられた顧客情報、課題、ソリューション要件を基に、課題の当事者にとって分かりやすいように補足し、提案書のPowerPointスライド2枚目のコンテンツをFunction Callingを用いて生成してください。"},
            {"role": "user", "content": f"顧客情報: {customer_info}\n課題: {issues}\nソリューション要件: {solution_requirements}"},
        ],
        functions=functions,
        function_call={"name": "summarize_issues_and_solutions"}
    )

    arguments = response.choices[0].message.function_call.arguments
    content = json.loads(arguments)

    return {
        "current_issues": content["current_issues"],
        "solution_direction": content["solution_direction"],
        "message_line": content["message_line"],
    }

def create_proposal_powerpoint(project_id: int, customer_info: str, bpmn_xml: str, issues: str, solution_requirements: str): # issues, solution_requirements を引数に追加
    """プロジェクトID、顧客情報、BPMN XML を基に提案書PowerPointを作成する。"""
    prs = Presentation(TEMPLATE_PPTX_PATH)

    # 1ページ目の処理: タイトルを "{customerInfo}様向け提案書" に変更
    slide1 = prs.slides[0]
    title_shape = slide1.shapes.title
    if title_shape:
        title_shape.text_frame.text = f"{customer_info}様向け提案書"

    # 2ページ目の処理: LLM で生成したテキストコンテンツを埋め込む
    slide2_content = generate_slide2_content(customer_info, issues, solution_requirements) # LLM でコンテンツを生成
    slide2 = prs.slides[1] # 2枚目のスライド

    def get_shape_by_name(slide, shape_name): # ヘルパー関数 get_shape_by_name を定義
        """スライドとShape名からShapeオブジェクトを取得するヘルパー関数"""
        for shape in slide.shapes:
            if shape.name == shape_name:
                return shape
        return None

    # Shapeオブジェクトを名前で取得
    issues_shape = get_shape_by_name(slide2, "正方形/長方形 1") # ヘルパー関数 get_shape_by_name を使用
    solution_direction_shape = get_shape_by_name(slide2, "正方形/長方形 5") # ヘルパー関数 get_shape_by_name を使用
    message_line_placeholder = slide2.shapes.placeholders[1] # 1つ目のプレースホルダー (メッセージライン用)

    # テキストコンテンツを Shape に設定
    if issues_shape:
        text_frame = issues_shape.text_frame
        text_frame.text = slide2_content["current_issues"]
        text_frame.word_wrap = True # 図形内でテキストを折り返す設定
        text_frame.vertical_anchor = MSO_ANCHOR.TOP # 上揃え
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT # 左揃え
            paragraph.font.name = "Meiryo UI" # Meiryo UI フォント

    if solution_direction_shape:
        text_frame = solution_direction_shape.text_frame
        text_frame.text = slide2_content["solution_direction"]
        text_frame.word_wrap = True # 図形内でテキストを折り返す設定
        text_frame.vertical_anchor = MSO_ANCHOR.TOP # 上揃え
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT # 左揃え
            paragraph.font.name = "Meiryo UI" # Meiryo UI フォント
    if message_line_placeholder:
        message_line_placeholder.text_frame.text = slide2_content["message_line"]

    # 3ページ目の処理: 業務フロー図を埋め込む (projects.py の create_powerpoint_file 関数を転用) 
    if bpmn_xml:
        slide3 = prs.slides[2] # 3枚目のスライド
        shapes_data, connectors = parse_bpmn_xml_content(bpmn_xml)
        if shapes_data:
            # BPMN 図形描画処理 (projects.py の create_powerpoint_file から転用)
            y_offset = Inches(1)
            shapes_on_slide3 = {} # slide3 に描画した図形を格納する辞書
            for key, data in shapes_data.items():
                shape_type_str = data["type"]
                shape_type = None

                if shape_type_str == "startEvent" or shape_type_str == "endEvent":
                    shape_type = MSO_SHAPE.OVAL
                elif shape_type_str == "task":
                    shape_type = MSO_SHAPE.RECTANGLE
                elif shape_type_str == "exclusiveGateway":
                    shape_type = MSO_SHAPE.DIAMOND
                else:
                    continue

                if data["type"] == "exclusiveGateway" and "outgoing" in data and data["outgoing"]:
                    first_outgoing_task_id = data["outgoing"][0]
                    if first_outgoing_task_id in shapes_data:
                        gateway_y = data["y"]
                        gateway_height = data["height"]
                        data["y"] = gateway_y + gateway_height / 2 + 0.5

                shape = slide3.shapes.add_shape( # slide3 に図形を追加
                    shape_type,
                    Inches(data["x"]), Inches(data["y"]) + y_offset, Inches(data["width"]), Inches(data["height"])
                )

                text_frame = shape.text_frame
                text_frame.text = data["name"]
                for paragraph in text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    paragraph.font.size = Inches(0.1)
                    paragraph.font.name = "Meiryo UI"
                    paragraph.font.color.rgb = RGBColor(0, 0, 0)
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
                shape.line.color.rgb = RGBColor(0, 0, 0)

                shapes_on_slide3[key] = shape # slide3 に描画した図形を shapes_on_slide3 に格納


            # コネクタ描画処理 (projects.py の create_powerpoint_file から転用、shapes_on_slide3 を使用)
            connectors_on_slide3 = connectors # BPMN図と同じコネクタ情報を使用
            for source, target in connectors_on_slide3:
                source_shape = shapes_on_slide3.get(source) # shapes_on_slide3 から図形を取得
                target_shape = shapes_on_slide3.get(target) # shapes_on_slide3 から図形を取得
                if source_shape and target_shape:
                    if source_shape.top != target_shape.top:
                        connector_type = MSO_CONNECTOR.ELBOW
                    else:
                        connector_type = MSO_CONNECTOR.STRAIGHT

                    connector = slide3.shapes.add_connector( # slide3 にコネクタを追加
                        connector_type, 0, 0, 0, 0
                    )

                    if connector_type == MSO_CONNECTOR.ELBOW:
                        if source_shape.left < target_shape.left:
                            connector.begin_connect(source_shape, 3)
                            connector.end_connect(target_shape, 1)
                        else:
                            connector.begin_connect(source_shape, 1)
                            connector.end_connect(target_shape, 3)
                    else:
                        connector.begin_connect(source_shape, 3)
                        connector.end_connect(target_shape, 1)

                    connector.line.width = Inches(0.02)
                    connector.line.color.rgb = RGBColor(0, 0, 0)

    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io

# PowerPoint ファイルのエンドポイント (提案書版)
@router.get("/{project_id}/proposal")
async def get_project_proposal(project_id: int = Path(..., gt=0)):
    df = read_projects()
    project_row = df[df['id'] == project_id]
    if project_row.empty:
        raise HTTPException(status_code=404, detail="Project not found")
    customer_name = project_row.iloc[0]['customer_name']
    bpmn_xml = project_row.iloc[0]['bpmn_xml']
    issues = project_row.iloc[0]['issues'] # 課題 (issues) を取得
    solution_requirements = project_row.iloc[0]['solution_requirements'] # ソリューション要件を取得

    pptx_file_path = None
    try:
        pptx_io = create_proposal_powerpoint(project_id, customer_name, bpmn_xml, issues, solution_requirements) # issues, solution_requirements を引数に追加
        return Response(
            content=pptx_io.getvalue(),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename=project_{project_id}_proposal.pptx"}
        )
    except HTTPException as e:
        raise e
    except Exception as e:
        error_message = f"提案書PowerPointオブジェクト生成エラー: {str(e)}"
        logging.error(error_message, exc_info=True)
        raise HTTPException(status_code=500, detail=error_message)
    finally:
        pass

@router.post("/")
async def create_proposal(proposal: ProposalCreate):
    solutions = read_csv("solutions.csv")
    solution = next((s for s in solutions if s["id"] == proposal.solution_id), None)
    
    if not solution:
        raise HTTPException(status_code=404, detail="Solution not found")

    generated_proposal = generate_proposal(
        proposal.customer_info,
        proposal.project_info,
        solution["name"] + ": " + solution["features"]
    )

    proposals = read_csv("proposals.csv")
    new_proposal = {
        "id": str(len(proposals) + 1),
        "customer_info": proposal.customer_info,
        "project_info": proposal.project_info,
        "solution_id": proposal.solution_id,
        "content": generated_proposal,
        "flagged": "false"
    }
    proposals.append(new_proposal)
    write_csv("proposals.csv", proposals)

    return new_proposal
